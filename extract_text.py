import sys
try:
    from pptx import Presentation
except ImportError:
    print("Please install python-pptx library first by running: pip install python-pptx")
    sys.exit(1)

def extract_text_from_pptx(filepath):
    print(f"Extracting text from: {filepath}")
    prs = Presentation(filepath)
    text_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text_content.append(f"\n--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
                
    output_filename = filepath.replace('.pptx', '_extracted.txt')
    with open(output_filename, 'w', encoding='utf-8') as f:
        f.write("\n".join(text_content))
    print(f"Extraction complete! Saved to {output_filename}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_text.py <path_to_pptx_file>")
    else:
        extract_text_from_pptx(sys.argv[1])
